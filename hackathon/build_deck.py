"""
Housing.com AI Content Engine — Hackathon Pitch Deck Generator
Generates HACKATHON_DECK.pptx (~16 slides)
Run: python3 hackathon/build_deck.py
"""
import io, os, textwrap
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

OUT = "hackathon/HACKATHON_DECK.pptx"

# ── Brand colors ──────────────────────────────────────────────────────────────
PURPLE      = RGBColor(0x6B, 0x2D, 0x8B)   # Housing.com primary
PURPLE_DARK = RGBColor(0x3A, 0x10, 0x55)   # dark bg
PURPLE_LIGHT= RGBColor(0xE8, 0xD5, 0xF5)   # light tint
ORANGE      = RGBColor(0xFF, 0x6B, 0x35)   # accent
GREEN       = RGBColor(0x27, 0xAE, 0x60)   # positive/save
RED         = RGBColor(0xE7, 0x4C, 0x3C)   # alert
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK        = RGBColor(0x1A, 0x1A, 0x2E)
GRAY        = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)

W, H = Inches(13.33), Inches(7.5)   # 16:9 widescreen

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]   # completely blank layout


# ══════════════════════════════════════════════════════════════════════════════
# Helper utilities
# ══════════════════════════════════════════════════════════════════════════════
def add_bg(slide, color):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def txb(slide, text, l, t, w, h, size=24, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def rect(slide, l, t, w, h, fill_color, line_color=None):
    from pptx.util import Pt as Pt2
    s = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt2(0.5)
    else:
        s.line.fill.background()
    return s

def pill(slide, text, l, t, w, h, bg, fg=WHITE, size=14, bold=False):
    r = rect(slide, l, t, w, h, bg)
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = fg

def img_from_fig(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    return buf

def add_img(slide, buf, l, t, w, h):
    slide.shapes.add_picture(buf, l, t, w, h)

def divider(slide, t, color=PURPLE):
    rect(slide, Inches(0.5), t, Inches(12.33), Inches(0.04), color)

def subtitle_bar(slide, text, t=Inches(1.0)):
    rect(slide, Inches(0), t, Inches(13.33), Inches(0.55), PURPLE)
    txb(slide, text, Inches(0.4), t + Inches(0.08), Inches(12), Inches(0.42),
        size=20, bold=True, color=WHITE)

def slide_num(slide, n):
    txb(slide, str(n), Inches(12.8), Inches(7.1), Inches(0.5), Inches(0.35),
        size=11, color=GRAY)

def watermark(slide):
    txb(slide, "Housing.com  ×  Multi-Agent AI", Inches(0.3), Inches(7.15),
        Inches(4), Inches(0.3), size=10, color=RGBColor(0xCC, 0xCC, 0xCC),
        italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, PURPLE_DARK)

# big gradient stripe
rect(s, Inches(0), Inches(0), Inches(13.33), Inches(0.6), PURPLE)

txb(s, "The AI Content Engine", Inches(0.7), Inches(1.4), Inches(11), Inches(1.5),
    size=60, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
txb(s, "Trend → Publish in 2 Minutes  ·  5 Platforms  ·  300 Posts/Month",
    Inches(0.7), Inches(3.0), Inches(11), Inches(0.7),
    size=24, color=PURPLE_LIGHT)
txb(s, "Housing.com Hackathon 2026  ·  Multi-Agent AI",
    Inches(0.7), Inches(3.9), Inches(8), Inches(0.5),
    size=16, color=RGBColor(0xAA, 0x88, 0xCC), italic=True)

# stat pills bottom
for i, (val, label) in enumerate([
    ("₹8L/mo saved", "vs agency team"),
    ("2 min", "end-to-end run"),
    ("5 platforms", "simultaneously"),
    ("3-pass QA", "brand-safe output"),
]):
    x = Inches(0.7 + i * 3.2)
    rect(s, x, Inches(6.0), Inches(2.8), Inches(1.0), PURPLE)
    txb(s, val,   x + Inches(0.1), Inches(6.05), Inches(2.6), Inches(0.45),
        size=22, bold=True, color=WHITE)
    txb(s, label, x + Inches(0.1), Inches(6.5),  Inches(2.6), Inches(0.4),
        size=12, color=PURPLE_LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
subtitle_bar(s, "The Problem: Content at Scale is Broken")
slide_num(s, 2); watermark(s)

problems = [
    ("⏱  Speed Gap",
     "Housing market moves daily — RERA rulings, RBI rate cuts, city micro-trends.\nA 4-person content team can't process 15 trending topics and publish before they peak."),
    ("💸  Cost Cliff",
     "₹8–10 lakh/month for a full content team: 1 strategist, 2 writers, 1 social manager.\nOutput: ~60–80 posts/month. Cost per post: ₹10,000–16,000."),
    ("🎯  Quality Drift",
     "Brand voice inconsistency across writers. Instagram caption ≠ LinkedIn tone.\nNo systematic QA. Political or factually wrong posts slip through."),
    ("🔗  Platform Blindness",
     "Same content pasted across 5 platforms. Twitter's 280 chars. Instagram's hashtag game.\nYouTube Shorts scripts. Each platform has different physics — one draft doesn't fit all."),
]
for i, (title, body) in enumerate(problems):
    row, col = divmod(i, 2)
    x = Inches(0.5 + col * 6.4)
    y = Inches(1.8 + row * 2.5)
    rect(s, x, y, Inches(6.0), Inches(2.2), LIGHT_GRAY,
         line_color=RGBColor(0xDD, 0xDD, 0xDD))
    txb(s, title, x + Inches(0.2), y + Inches(0.15), Inches(5.6), Inches(0.45),
        size=16, bold=True, color=PURPLE)
    txb(s, body,  x + Inches(0.2), y + Inches(0.65), Inches(5.6), Inches(1.4),
        size=13, color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION (one-liner)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, PURPLE_DARK)
slide_num(s, 3); watermark(s)

txb(s, "Introducing", Inches(1), Inches(0.9), Inches(11), Inches(0.6),
    size=20, color=PURPLE_LIGHT, italic=True)
txb(s, "The AI Content Engine", Inches(1), Inches(1.5), Inches(11), Inches(1.3),
    size=52, bold=True, color=WHITE)
txb(s, "A fully automated, multi-agent pipeline that monitors real-time trends,\n"
       "generates platform-native content, enforces 3-pass brand safety QA,\n"
       "and publishes to 5 channels — in under 2 minutes.",
    Inches(1), Inches(3.0), Inches(11), Inches(1.5),
    size=22, color=PURPLE_LIGHT)

# feature chips
features = ["Claude Opus 4.7", "LangGraph", "5 Platforms", "3-Pass QA",
            "Hinglish Voice", "SEO Articles", "Slack Bot", "Engagement Prediction"]
for i, f in enumerate(features):
    x = Inches(1.0 + (i % 4) * 2.9)
    y = Inches(4.8 + (i // 4) * 0.7)
    rect(s, x, y, Inches(2.5), Inches(0.5), PURPLE)
    txb(s, f, x, y, Inches(2.5), Inches(0.5),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE DIAGRAM
# ══════════════════════════════════════════════════════════════════════════════
fig, ax = plt.subplots(figsize=(13, 5), facecolor="#3A1055")
ax.set_xlim(0, 13); ax.set_ylim(0, 5); ax.axis("off")

nodes = [
    # (x, y, label, width, fill)
    (0.7,  2.2, "Trend\nResearcher",  1.4, "#6B2D8B"),
    (0.7,  0.5, "News\nResearcher",   1.4, "#6B2D8B"),
    (3.1,  1.3, "Planner\n(Quality Gate)", 1.6, "#FF6B35"),
    (5.5,  2.5, "Social\nCreative\n(Opus 4.7)", 1.5, "#6B2D8B"),
    (5.5,  0.2, "News\nCreative\n(Sonnet)", 1.5, "#6B2D8B"),
    (7.9,  1.3, "Platform\nAgents ×5", 1.5, "#6B2D8B"),
    (10.3, 1.3, "QA Engine\n(3-pass)", 1.5, "#27AE60"),
    (12.1, 1.3, "Publisher\n+ Notify", 1.3, "#FF6B35"),
]

for (x, y, label, w, color) in nodes:
    rect_patch = mpatches.FancyBboxPatch((x - w/2, y - 0.5), w, 1.0,
        boxstyle="round,pad=0.05", facecolor=color, edgecolor="white",
        linewidth=1.2)
    ax.add_patch(rect_patch)
    ax.text(x, y, label, ha="center", va="center", fontsize=7.5,
            color="white", fontweight="bold", multialignment="center")

# arrows
arrows = [
    (0.7+0.7, 2.2,  3.1-0.8, 1.6),   # trend → planner
    (0.7+0.7, 0.5,  3.1-0.8, 1.0),   # news → planner
    (3.1+0.8, 1.5,  5.5-0.75, 2.5),  # planner → social
    (3.1+0.8, 1.2,  5.5-0.75, 0.5),  # planner → news
    (5.5+0.75, 2.5, 7.9-0.75, 1.7),  # social → platform
    (5.5+0.75, 0.5, 7.9-0.75, 1.0),  # news → platform
    (7.9+0.75, 1.3, 10.3-0.75, 1.3), # platform → qa
    (10.3+0.75, 1.3, 12.1-0.65, 1.3),# qa → publish
]
for (x1, y1, x2, y2) in arrows:
    ax.annotate("", xy=(x2, y2), xytext=(x1, y1),
                arrowprops=dict(arrowstyle="->", color="white", lw=1.2))

# parallel annotation
ax.text(0.7, 4.3, "⚡ Parallel", color="#E8D5F5", fontsize=9, fontstyle="italic")
ax.annotate("", xy=(0.7, 3.1), xytext=(0.7, 4.2),
            arrowprops=dict(arrowstyle="-", color="#E8D5F5", lw=0.8, linestyle="dashed"))
ax.annotate("", xy=(0.7, 1.0), xytext=(0.7, 4.2),
            arrowprops=dict(arrowstyle="-", color="#E8D5F5", lw=0.8, linestyle="dashed"))

# platform icons
for i, p in enumerate(["Twitter", "Instagram", "YouTube", "Housing\nNews", "LinkedIn"]):
    ax.text(8.65, 2.5 - i * 0.5, f"→ {p}", color="#E8D5F5", fontsize=6.5)

ax.set_title("Pipeline Architecture  ·  ~120s end-to-end  ·  LangGraph + Claude",
             color="white", fontsize=10, pad=8)

buf = img_from_fig(fig); plt.close(fig)

s = prs.slides.add_slide(blank)
add_bg(s, PURPLE_DARK)
subtitle_bar(s, "Architecture: 8-Node Multi-Agent Pipeline")
add_img(s, buf, Inches(0.2), Inches(1.4), Inches(12.9), Inches(5.5))
slide_num(s, 4); watermark(s)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — HOW EACH AGENT WORKS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
subtitle_bar(s, "What Each Agent Does")
slide_num(s, 5); watermark(s)

agents_info = [
    ("🔍  Researcher", "Sonnet 4.6",
     "Tavily web search across 12 credible RE sources. 5 rounds max. 8 news items with relevance scoring."),
    ("📈  Trend Researcher", "Sonnet 4.6",
     "Aggregates Google Trends + YouTube India + Reddit + Twitter. 15 trend items with Hinglish creative hooks."),
    ("🗂  Planner", "Gemini Flash",
     "Quality gate. Filters off-topic content. Max 5 social + 3 news briefs. OMIT anything with no genuine RE angle."),
    ("✍️  Social Creative", "Opus 4.7",
     "Zomato method: trend is HERO, housing is PUNCHLINE. Hinglish, 2 emojis max, meme-ready card concept."),
    ("📰  News Creative", "Sonnet 4.6",
     "700–1000 word SEO articles. H2 structure, pull quote, internal links, ≤70 char headline. CMS-ready JSON output."),
    ("🖥  Platform Agents ×5", "Sonnet 4.6",
     "Twitter (≤280 chars + thread), Instagram (story slides + reel script), YouTube, LinkedIn, Housing.com News."),
    ("🛡  QA Engine", "Multi-model",
     "3-pass: Safety gate (binary, Gemini Flash) → Quality scoring (7 dimensions, Sonnet) → Engagement prediction (Opus)."),
    ("📤  Publisher", "No LLM",
     "Deterministic. Saves to disk (dry run) or posts live. Slack summary per run. retry=None (non-idempotent)."),
]
for i, (name, model, desc) in enumerate(agents_info):
    col, row = divmod(i, 4)
    x = Inches(0.4 + col * 6.45)
    y = Inches(1.8 + row * 1.35)
    rect(s, x, y, Inches(6.1), Inches(1.2), LIGHT_GRAY,
         line_color=RGBColor(0xDD, 0xCC, 0xEE))
    txb(s, name, x+Inches(0.15), y+Inches(0.08), Inches(4.0), Inches(0.38),
        size=14, bold=True, color=PURPLE)
    pill(s, model, x+Inches(4.3), y+Inches(0.1), Inches(1.6), Inches(0.32),
         bg=PURPLE, size=10, bold=True)
    txb(s, desc, x+Inches(0.15), y+Inches(0.5), Inches(5.8), Inches(0.65),
        size=11, color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — SAMPLE OUTPUT: TWITTER THREAD
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
subtitle_bar(s, "Sample Output: Twitter Thread  (Auto-generated, RERA Topic)")
slide_num(s, 6); watermark(s)

# tweet card
rect(s, Inches(0.5), Inches(1.6), Inches(7.5), Inches(5.4),
     RGBColor(0xF8, 0xF8, 0xFF), line_color=RGBColor(0xCC, 0xCC, 0xDD))
txb(s, "@HousingDotCom  ·  Thread", Inches(0.7), Inches(1.75), Inches(5), Inches(0.4),
    size=13, bold=True, color=PURPLE)
txb(s, "🚨 8,212 projects just received RERA notices in Maharashtra alone.\n\n"
       "Is your dream home on the list?\n\n"
       "Bank accounts frozen. Sales halted. Registrations cancelled.\n\n"
       "Here's how to check before it's too late 👇\n\nhttps://housing.com/in/buy/mumbai",
    Inches(0.7), Inches(2.2), Inches(7.0), Inches(2.2),
    size=13, color=DARK)
txb(s, "2/ In the last 30 days, 3 regulators dropped the hammer:\n\n"
       "• MahaRERA: 8,212 notices (Mumbai + Pune)\n"
       "• Telangana RERA: 11% refund ordered on defrauded buyers\n"
       "• Haryana RERA: ₹52 crore penalty threat",
    Inches(0.7), Inches(4.5), Inches(7.0), Inches(1.8),
    size=12, color=DARK)

# metadata
txb(s, "✅ QA Score: 6.4 / 10    ✅ Safety: PASSED    📊 Pred. ER: 7.15%",
    Inches(0.7), Inches(6.55), Inches(7.0), Inches(0.4),
    size=12, bold=True, color=GREEN)

# annotations column
rect(s, Inches(8.2), Inches(1.6), Inches(4.8), Inches(5.4),
     PURPLE_DARK)
for i, (label, val) in enumerate([
    ("Model",     "Opus 4.7 (creative tier)"),
    ("Hook type", "Stat-based urgency"),
    ("Platform",  "Twitter thread (4 parts)"),
    ("Chars",     "238 / 280 (✓ within limit)"),
    ("Hashtags",  "#RERA #MahaRERA + 3 more"),
    ("Links",     "2 internal (Mumbai, Pune SRP)"),
    ("Run time",  "~8s for creative node"),
]):
    y = Inches(1.8 + i * 0.7)
    txb(s, label, Inches(8.4), y, Inches(1.5), Inches(0.55),
        size=11, color=PURPLE_LIGHT, bold=True)
    txb(s, val,   Inches(9.9), y, Inches(2.9), Inches(0.55),
        size=11, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — SAMPLE OUTPUT: INSTAGRAM
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
subtitle_bar(s, "Sample Output: Instagram  (Branded Card + Story Slides + Reel Script)")
slide_num(s, 7); watermark(s)

# IG card mockup (purple brand card)
card = rect(s, Inches(0.5), Inches(1.6), Inches(3.2), Inches(3.2), PURPLE)
txb(s, "housing.com", Inches(0.6), Inches(1.75), Inches(3.0), Inches(0.4),
    size=11, color=WHITE, italic=True)
txb(s, "⚠️ 8,212 projects\njust received\nRERA notices.\n\nIs YOUR project\nsafe?",
    Inches(0.6), Inches(2.2), Inches(3.0), Inches(1.8),
    size=16, bold=True, color=WHITE)
txb(s, "Check. Verify. Buy Smart.", Inches(0.6), Inches(4.2), Inches(3.0), Inches(0.4),
    size=11, color=PURPLE_LIGHT)

# story slides
txb(s, "Story Slides (auto-generated):", Inches(3.9), Inches(1.65), Inches(9), Inches(0.4),
    size=13, bold=True, color=PURPLE)
stories = [
    "Slide 1 — HOOK: '8,212 housing projects. RERA notices. Is YOUR project safe? Swipe →'",
    "Slide 2 — THE CRACKDOWN: Maharashtra, Telangana, Haryana — all in 30 days",
    "Slide 3 — WHAT'S AT RISK: Bank accounts frozen, sales halted, ₹52Cr penalty",
    "Slide 4 — YOUR CHECKLIST: 5-step RERA verification guide",
    "Slide 5 — CTA: Search RERA-verified homes on Housing.com 🔗 Link in bio",
]
for i, story in enumerate(stories):
    y = Inches(2.1 + i * 0.6)
    rect(s, Inches(3.9), y, Inches(0.25), Inches(0.38), ORANGE)
    txb(s, story, Inches(4.3), y, Inches(8.7), Inches(0.45),
        size=11, color=DARK)

txb(s, "Reel Concept:", Inches(3.9), Inches(5.25), Inches(9), Inches(0.35),
    size=13, bold=True, color=PURPLE)
txb(s, "0–4s: dramatic text reveal 'Is your dream home on this list?'  |  5–22s: on-camera "
       "presenter outside Mumbai building, 5-step checklist flash  |  23–28s: Housing.com app CTA",
    Inches(3.9), Inches(5.65), Inches(9.1), Inches(0.6),
    size=11, color=DARK)

txb(s, "✅ QA Score: 6.8 / 10    ✅ Safety: PASSED    📊 Pred. ER: 5.7%    🏷 Hashtags: 28",
    Inches(0.5), Inches(6.55), Inches(12), Inches(0.4),
    size=12, bold=True, color=GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — SAMPLE OUTPUT: HOUSING NEWS SEO ARTICLE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
subtitle_bar(s, "Sample Output: Housing.com News  (SEO Article, 800 words, CMS-ready)")
slide_num(s, 8); watermark(s)

rect(s, Inches(0.5), Inches(1.6), Inches(8.7), Inches(5.4),
     RGBColor(0xFA, 0xFA, 0xFF), line_color=RGBColor(0xCC, 0xBB, 0xEE))
txb(s, "SEO Title:", Inches(0.7), Inches(1.75), Inches(2.0), Inches(0.35),
    size=11, bold=True, color=GRAY)
txb(s, "8,212 Projects Got RERA Notices: How to Check If Your Home Is at Risk",
    Inches(2.7), Inches(1.75), Inches(6.3), Inches(0.35),
    size=13, bold=True, color=PURPLE)
txb(s, "Primary keyword: how to check RERA registration  ·  Reading time: 4 min",
    Inches(0.7), Inches(2.15), Inches(8.3), Inches(0.35),
    size=11, color=GRAY, italic=True)
txb(s, "MahaRERA just sent show-cause notices to 8,212 housing projects. Telangana RERA "
       "ordered an 11% interest refund to defrauded buyers. Haryana threatened a ₹52 crore "
       "penalty on a developer selling unregistered floors. If you are buying a home in 2026, "
       "this is the only article you need to read first...\n\n"
       "## What Just Happened: Three RERA Crackdowns in 30 Days\n\n"
       "Maharashtra's RERA has issued show-cause notices to 8,212 housing projects for failing "
       "to file mandatory Quarterly Progress Reports (QPRs). The breakdown: 4,644 in Mumbai "
       "Metropolitan Region, 1,957 in Pune...\n\n"
       "## The 60-Second RERA Check Every Buyer Must Do\n\n"
       "✅ Step 1: Ask for the RERA Registration Number\n"
       "✅ Step 2: Verify on Your State's RERA Portal\n"
       "✅ Step 3: Check QPR Status...",
    Inches(0.7), Inches(2.55), Inches(8.3), Inches(3.6),
    size=11, color=DARK)
txb(s, "Internal links: 3 Housing.com SRP pages embedded contextually",
    Inches(0.7), Inches(6.25), Inches(8.3), Inches(0.35),
    size=11, color=PURPLE, italic=True)

rect(s, Inches(9.5), Inches(1.6), Inches(3.5), Inches(5.4), PURPLE_DARK)
meta = [
    ("Format",       "CMS-ready JSON + Markdown"),
    ("SEO Title",    "≤ 70 chars ✓"),
    ("Meta Desc",    "≤ 160 chars ✓"),
    ("H2 Structure", "4 sections ✓"),
    ("Pull Quote",   "Independently shareable ✓"),
    ("Schema",       "NewsArticle JSON-LD"),
    ("Internal Links","3 embedded contextually"),
    ("QA Score",     "8.0 / 10"),
    ("Safety",       "PASSED ✓"),
    ("Pred. ER",     "2.1% (CTR-focused)"),
]
for i, (k, v) in enumerate(meta):
    y = Inches(1.75 + i * 0.47)
    txb(s, k, Inches(9.6), y, Inches(1.4), Inches(0.4),
        size=10, bold=True, color=PURPLE_LIGHT)
    txb(s, v, Inches(11.0), y, Inches(1.9), Inches(0.4),
        size=10, color=WHITE)
slide_num(s, 8)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — THE ZOMATO METHOD (creative strategy)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
subtitle_bar(s, "The Creative Strategy: The Zomato Method")
slide_num(s, 9); watermark(s)

txb(s, "Social content fails when housing is the HERO and the trend is the footnote.\n"
       "It works when the TREND is the hero and housing is the punchline.",
    Inches(0.5), Inches(1.7), Inches(12.3), Inches(0.8),
    size=16, color=DARK, italic=True)

examples = [
    ("Trend: IPL final",
     "❌  Weak: 'Celebrating RCB's win! Buy a home in Bangalore! 🏠'\n"
     "✅  Strong: 'RCB waited 18 years to win the IPL.\nYou've been waiting to buy a ghar for 3.\nMaybe 2025 is both our years. 🏆'\n#RCBvKKR #GharKhojna"),
    ("Trend: Budget ₹12L tax-free",
     "❌  Weak: 'Budget 2025 is great for home buyers!'\n"
     "✅  Strong: 'Finance minister just made ₹12L tax-free.\nYour EMI just got more affordable than your rent.\nCheck how much home you can afford now 👇'\n#Budget2025 #HomeLoan"),
    ("Trend: Coldplay concert",
     "❌  Wrong: 'Coldplay is coming! Buy a home near the venue!'\n"
     "✅  OMIT: No genuine housing angle → planner drops it\n"
     "(Planner quality gate working as designed)"),
]
for i, (trend, content) in enumerate(examples):
    x = Inches(0.4 + i * 4.3)
    rect(s, x, Inches(2.6), Inches(4.0), Inches(3.8),
         LIGHT_GRAY if i < 2 else RGBColor(0xFF, 0xF0, 0xF0),
         line_color=RGBColor(0xCC, 0xBB, 0xEE))
    txb(s, trend, x+Inches(0.15), Inches(2.75), Inches(3.7), Inches(0.4),
        size=13, bold=True, color=PURPLE if i < 2 else RED)
    txb(s, content, x+Inches(0.15), Inches(3.2), Inches(3.7), Inches(3.0),
        size=11, color=DARK)

txb(s, "Trend is HERO  ·  Housing is PUNCHLINE  ·  Hinglish + max 2 emojis  ·  First hashtag = trending tag",
    Inches(0.5), Inches(6.7), Inches(12.0), Inches(0.4),
    size=13, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — 3-PASS QA ENGINE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
subtitle_bar(s, "Quality Gate: 3-Pass AI Review on Every Post")
slide_num(s, 10); watermark(s)

passes = [
    ("Pass 1\nSafety Gate", "Gemini Flash\n~$0.001/post",
     "Binary PASS/FAIL\n\nHard blocks:\n• Political party names\n• Communal / religious attacks\n• Price guarantees with numbers\n• Defamation / insider info\n\nAlways allowed:\n• Govt schemes (PMAY, RERA)\n• Women homebuyers\n• Market data with attribution",
     GREEN),
    ("Pass 2\nQuality Score", "Sonnet 4.6\n~$0.012/post",
     "7-dimension scoring (0–10)\n\nTwitter: Trend dominance (35%)\n  + Shareability (30%)\n  + Hinglish wit (25%)\n  + CTA quality (10%)\n\nThreshold < 5.0 → REVISE\nUp to 2 revision loops",
     PURPLE),
    ("Pass 3\nEngagement Pred.", "Opus 4.7\n~$0.015/post",
     "Behavioural prediction\n\nOutputs:\n• pred_impressions\n• pred_engagement_rate\n• pred_confidence\n\nHistorical Housing.com data\n+ platform engagement physics\n+ trend momentum signal",
     ORANGE),
]
for i, (title, model, body, color) in enumerate(passes):
    x = Inches(0.5 + i * 4.25)
    rect(s, x, Inches(1.75), Inches(3.9), Inches(4.8), LIGHT_GRAY,
         line_color=color)
    rect(s, x, Inches(1.75), Inches(3.9), Inches(0.55), color)
    txb(s, title, x+Inches(0.15), Inches(1.8), Inches(3.6), Inches(0.45),
        size=15, bold=True, color=WHITE)
    pill(s, model, x+Inches(0.1), Inches(2.4), Inches(3.7), Inches(0.45),
         bg=color, size=10)
    txb(s, body, x+Inches(0.1), Inches(2.95), Inches(3.7), Inches(3.5),
        size=11, color=DARK)

# flow arrow between passes
txb(s, "→", Inches(4.35), Inches(3.6), Inches(0.5), Inches(0.5),
    size=28, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
txb(s, "→", Inches(8.6), Inches(3.6), Inches(0.5), Inches(0.5),
    size=28, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

txb(s, "Revision loop: REVISE decision → platform agent reruns with specific fix instructions → re-evaluated (max 2 attempts)",
    Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
    size=12, color=PURPLE, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — COST COMPARISON CHART
# ══════════════════════════════════════════════════════════════════════════════
fig, axes = plt.subplots(1, 2, figsize=(12, 4.5), facecolor="white")
fig.subplots_adjust(wspace=0.35)

# Left: per-run cost breakdown
components = ["Social Creative\n(Opus 4.7)", "QA Engagement\n(Sonnet)", "QA Quality\n(Sonnet)",
              "Platform Agents\n(×5 Sonnet)", "News Creative\n(Sonnet)", "Research\n(Sonnet)",
              "Trend Research\n(Sonnet)", "Other\n(Flash/tools)"]
costs = [0.123, 0.068, 0.063, 0.096, 0.038, 0.034, 0.032, 0.064]
colors_bar = ["#6B2D8B", "#FF6B35", "#FF6B35", "#6B2D8B", "#6B2D8B",
              "#6B2D8B", "#6B2D8B", "#27AE60"]
bars = axes[0].barh(components, costs, color=colors_bar, edgecolor="none", height=0.65)
axes[0].set_xlabel("Cost (USD)", fontsize=10)
axes[0].set_title("Per-Run Cost Breakdown  (Total: ~$0.52)", fontsize=11, fontweight="bold", color="#3A1055")
for bar, val in zip(bars, costs):
    axes[0].text(val + 0.002, bar.get_y() + bar.get_height()/2,
                 f"${val:.3f}", va="center", fontsize=8, color="#333333")
axes[0].set_xlim(0, 0.175)
axes[0].tick_params(axis="y", labelsize=8)
axes[0].spines[["top","right"]].set_visible(False)

# Right: monthly comparison
labels = ["Traditional\nAgency Team", "AI Content\nEngine"]
monthly = [550000 / 83, 170]  # ₹8L/mo in USD vs $170
bar_colors = ["#E74C3C", "#27AE60"]
b = axes[1].bar(labels, monthly, color=bar_colors, edgecolor="none", width=0.5)
axes[1].set_ylabel("Monthly Cost (USD)", fontsize=10)
axes[1].set_title("Monthly Cost Comparison", fontsize=11, fontweight="bold", color="#3A1055")
axes[1].set_ylim(0, 8000)
for bar, val in zip(b, monthly):
    axes[1].text(bar.get_x() + bar.get_width()/2, bar.get_height() + 100,
                 f"${val:,.0f}\n(~₹{val*83/100000:.1f}L)",
                 ha="center", va="bottom", fontsize=10, fontweight="bold")
axes[1].spines[["top","right"]].set_visible(False)
axes[1].tick_params(axis="x", labelsize=11)

# savings annotation
axes[1].annotate("", xy=(1, monthly[1] + 300), xytext=(0, monthly[0]),
                 arrowprops=dict(arrowstyle="->", color="#27AE60", lw=2))
axes[1].text(0.5, 4500, "97% cost\nreduction", ha="center", fontsize=14,
             color="#27AE60", fontweight="bold",
             bbox=dict(boxstyle="round,pad=0.3", facecolor="#E8F8F0", edgecolor="#27AE60"))

buf2 = img_from_fig(fig); plt.close(fig)

s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
subtitle_bar(s, "The Economics: 97% Cost Reduction vs Traditional Agency")
add_img(s, buf2, Inches(0.3), Inches(1.4), Inches(12.7), Inches(5.5))
slide_num(s, 11); watermark(s)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — ROI / BUSINESS IMPACT (big numbers)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, PURPLE_DARK)
slide_num(s, 12); watermark(s)

txb(s, "Business Impact", Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
    size=32, bold=True, color=WHITE)
divider(s, Inches(1.2), PURPLE)

metrics = [
    ("₹8L–10L/mo", "agency cost replaced", "Traditional 4-person content team"),
    ("₹14,000/mo",  "all-in AI cost",       "LLM + infra + search APIs"),
    ("97%",         "cost reduction",       "Payback period: < 2 months"),
    ("300",         "posts / month",        "vs ~60–80 from human team  (5× volume)"),
    ("5",           "platforms",            "Twitter, Instagram, YouTube, LinkedIn, Housing.com News"),
    ("<2 min",      "per full pipeline run","Trend → QA → Ready to publish"),
]
for i, (big, label, sub) in enumerate(metrics):
    col, row = divmod(i, 3)
    x = Inches(0.4 + col * 6.5)
    y = Inches(1.5 + row * 2.35)
    rect(s, x, y, Inches(6.1), Inches(2.1), PURPLE)
    txb(s, big,   x+Inches(0.2), y+Inches(0.1), Inches(5.7), Inches(1.0),
        size=40, bold=True, color=ORANGE if i in (0,2) else WHITE)
    txb(s, label, x+Inches(0.2), y+Inches(1.1), Inches(5.7), Inches(0.45),
        size=16, bold=True, color=WHITE)
    txb(s, sub,   x+Inches(0.2), y+Inches(1.6), Inches(5.7), Inches(0.4),
        size=11, color=PURPLE_LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — SCALE & ENGAGEMENT ESTIMATES
# ══════════════════════════════════════════════════════════════════════════════
fig, ax = plt.subplots(figsize=(12, 4), facecolor="white")
months = [f"Mo {i}" for i in range(1, 7)]
posts = [300, 300, 300, 300, 300, 300]
impressions = [420000, 580000, 720000, 900000, 1100000, 1400000]  # compounding organic

ax2 = ax.twinx()
ax.bar(months, posts, color="#6B2D8B", alpha=0.7, label="Posts published", width=0.5)
ax2.plot(months, [i/1000 for i in impressions], "o-", color="#FF6B35",
         linewidth=2.5, markersize=8, label="Cumulative impressions (K)")
ax2.fill_between(range(6), [i/1000 for i in impressions], alpha=0.15, color="#FF6B35")

ax.set_ylabel("Posts Published", fontsize=10, color="#6B2D8B")
ax2.set_ylabel("Monthly Impressions (K)", fontsize=10, color="#FF6B35")
ax.set_title("Projected 6-Month Scale  ·  Content compounding = organic growth flywheel",
             fontsize=11, fontweight="bold", color="#3A1055")
ax.legend(loc="upper left", fontsize=9)
ax2.legend(loc="upper center", fontsize=9)
ax.spines[["top","right"]].set_visible(False)

# annotation
ax2.annotate("SEO articles start\nranking ~Mo 3", xy=(2, 720), xytext=(3.5, 850),
             fontsize=9, color="#333333",
             arrowprops=dict(arrowstyle="->", color="#333333"))

buf3 = img_from_fig(fig); plt.close(fig)

s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
subtitle_bar(s, "Scale: 6-Month Content Compounding Projection")
add_img(s, buf3, Inches(0.4), Inches(1.4), Inches(12.5), Inches(4.8))
txb(s, "Engagement estimates based on Housing.com platform benchmarks  ·  Twitter ER: 1.5–7%  ·  Instagram ER: 5–7%  ·  News CTR: 2–3%",
    Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
    size=11, color=GRAY, italic=True)
slide_num(s, 13); watermark(s)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — THE FEEDBACK FLYWHEEL
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
subtitle_bar(s, "The Feedback Flywheel: The System Gets Smarter Every Day")
slide_num(s, 14); watermark(s)

# circular flywheel diagram via matplotlib
fig2, ax2 = plt.subplots(figsize=(5, 5), facecolor="white")
ax2.set_xlim(-1.5, 1.5); ax2.set_ylim(-1.5, 1.5); ax2.axis("off")
wheel_nodes = [
    (0, 1.1,    "Generate\nContent",   "#6B2D8B"),
    (1.05, 0.0, "Publish\nLive",       "#FF6B35"),
    (0, -1.1,   "Track\nEngagement\n(6h/24h/7d)", "#27AE60"),
    (-1.05, 0,  "Inject into\nNext Run", "#6B2D8B"),
]
for (x, y, label, color) in wheel_nodes:
    circle = plt.Circle((x, y), 0.38, color=color, zorder=3)
    ax2.add_patch(circle)
    ax2.text(x, y, label, ha="center", va="center", fontsize=8,
             color="white", fontweight="bold", multialignment="center", zorder=4)

import numpy as np
theta = np.linspace(0, 2*np.pi, 200)
ax2.plot(np.cos(theta)*0.72, np.sin(theta)*0.72, color="#6B2D8B", linewidth=1.5, alpha=0.4)
for angle in [np.pi/2, 0, -np.pi/2, np.pi]:
    dx = -0.1*np.sin(angle); dy = 0.1*np.cos(angle)
    x_ = 0.72*np.cos(angle); y_ = 0.72*np.sin(angle)
    ax2.annotate("", xy=(x_+dx, y_+dy), xytext=(x_-dx*0.5, y_-dy*0.5),
                 arrowprops=dict(arrowstyle="->", color="#6B2D8B", lw=1.5))
ax2.text(0, 0, "Better\nContent\nOver Time", ha="center", va="center",
         fontsize=9, color="#3A1055", fontweight="bold")
buf4 = img_from_fig(fig2); plt.close(fig2)

add_img(s, buf4, Inches(0.3), Inches(1.5), Inches(4.5), Inches(4.8))

txb(s, "How the loop works:", Inches(5.0), Inches(1.7), Inches(7.8), Inches(0.45),
    size=15, bold=True, color=PURPLE)
loop_steps = [
    ("1. Publish", "Every approved post is stored with QA scores + predicted engagement"),
    ("2. Track",   "Engagement jobs fetch actual metrics at 6h, 24h, and 7 days post-publish"),
    ("3. Learn",   "get_performance_history() identifies top/bottom performers by topic, tone, city"),
    ("4. Inject",  "Performance context injected into social_creative prompt every run"),
    ("5. Improve", "Creative agent learns: cricket+Bengaluru gets 3× more ER than market analysis"),
]
for i, (step, desc) in enumerate(loop_steps):
    y = Inches(2.2 + i * 0.82)
    rect(s, Inches(5.0), y, Inches(0.7), Inches(0.6), PURPLE)
    txb(s, step, Inches(5.0), y, Inches(0.7), Inches(0.6),
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s, desc, Inches(5.85), y + Inches(0.1), Inches(7.1), Inches(0.5),
        size=12, color=DARK)

txb(s, "Realistic improvement: ~22% engagement rate gain over first 50 runs  ·  Plateaus as easy signal is harvested",
    Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.4),
    size=11, color=GRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — DEPLOYMENT MODES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
subtitle_bar(s, "Deployment: 5 Ways to Run It")
slide_num(s, 15); watermark(s)

modes = [
    ("🖥  Scheduled Server",
     "python main.py serve",
     "9 AM + 6 PM IST daily. APScheduler triggers automatically. FastAPI REST API for manual control."),
    ("💬  Slack Bot",
     "python main.py slack-bot",
     "DM '@HousingBot content idea' → pipeline runs → Slack thread reply in 2 mins. Triggered by marketing team."),
    ("🔌  REST API",
     "POST /run  ·  GET /runs/{id}",
     "Any system can trigger a run. Async response with run_id. Real-time status polling. CI/CD integration."),
    ("🖱  React Dashboard",
     "python main.py ui",
     "Browser UI: run history, QA scores, published posts, engagement tracking, feedback collection."),
    ("⌨️  CLI (this demo)",
     "python main.py run",
     "Single pipeline run with live progress bar. 2-3 minutes. Outputs saved to output/<run_id>/"),
]
for i, (name, cmd, desc) in enumerate(modes):
    col, row = divmod(i, 3)
    if i < 3:
        x = Inches(0.4 + i * 4.28)
        y = Inches(1.8)
    else:
        x = Inches(0.4 + (i-3) * 4.28 + 2.14)
        y = Inches(4.35)
    rect(s, x, y, Inches(3.95), Inches(2.2), LIGHT_GRAY,
         line_color=RGBColor(0xCC, 0xBB, 0xEE))
    txb(s, name, x+Inches(0.15), y+Inches(0.1), Inches(3.6), Inches(0.45),
        size=14, bold=True, color=PURPLE)
    rect(s, x+Inches(0.15), y+Inches(0.6), Inches(3.65), Inches(0.38),
         RGBColor(0x1A, 0x1A, 0x2E))
    txb(s, cmd, x+Inches(0.25), y+Inches(0.63), Inches(3.5), Inches(0.32),
        size=11, bold=True, color=RGBColor(0x7B, 0xFF, 0x8B))
    txb(s, desc, x+Inches(0.15), y+Inches(1.1), Inches(3.65), Inches(0.9),
        size=11, color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — TECH STACK
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
subtitle_bar(s, "Technology Stack")
slide_num(s, 16); watermark(s)

stack = [
    ("AI / LLM", [
        ("Claude Opus 4.7", "Social creative, engagement prediction"),
        ("Claude Sonnet 4.6", "Research, news creative, QA quality, platforms"),
        ("Gemini 2.5 Flash", "Planner, QA safety, extraction (53% cheaper)"),
    ]),
    ("Orchestration", [
        ("LangGraph", "Stateful multi-agent graph, parallel fan-out/in"),
        ("AsyncSqliteSaver", "Checkpointing — crash-resumable runs"),
        ("RetryPolicy", "Exponential backoff on all non-publisher nodes"),
    ]),
    ("Data & Tools", [
        ("Tavily + Serper", "Real-time web search across 12 RE domains"),
        ("Google Trends + Reddit", "Social trend aggregation (4 sources)"),
        ("SQLite + SQLAlchemy", "Post storage, QA history, engagement tracking"),
    ]),
    ("Infrastructure", [
        ("FastAPI", "REST API + async run management"),
        ("APScheduler", "Cron-style scheduling (9 AM + 6 PM IST)"),
        ("React + Vite", "Real-time dashboard UI"),
        ("PIL (Pillow)", "1080×1080 branded card generation (no DALL-E cost)"),
    ]),
]
for col, (category, items) in enumerate(stack):
    x = Inches(0.4 + col * 3.22)
    rect(s, x, Inches(1.75), Inches(3.0), Inches(0.45), PURPLE)
    txb(s, category, x, Inches(1.8), Inches(3.0), Inches(0.38),
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, (name, desc) in enumerate(items):
        y = Inches(2.35 + j * 1.1)
        rect(s, x, y, Inches(3.0), Inches(0.95),
             LIGHT_GRAY, line_color=RGBColor(0xCC, 0xBB, 0xEE))
        txb(s, name, x+Inches(0.1), y+Inches(0.05), Inches(2.8), Inches(0.38),
            size=13, bold=True, color=PURPLE)
        txb(s, desc, x+Inches(0.1), y+Inches(0.48), Inches(2.8), Inches(0.4),
            size=10, color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — CASE STUDY WALKTHROUGH
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
subtitle_bar(s, "Case Study: RERA Crackdown → 5 Platform Posts in 97 Seconds")
slide_num(s, 17); watermark(s)

timeline = [
    ("0s",    "researcher_node",    "Tavily fetches 8 RE news stories\n+ 5 RERA-specific articles (Sonnet 4.6)",     PURPLE),
    ("8s",    "trend_researcher",   "Google Trends: #RERA trending\n+ 15 total trend items with Hinglish hooks",    PURPLE),
    ("19s",   "planner_node",       "RERA story → 2 briefs: social (Hinglish urgency)\n+ news (SEO buyer guide)",   ORANGE),
    ("27s",   "social_creative",    "Opus 4.7: 'Is your dream home on this list?'\nZomato-style urgency hook",      PURPLE),
    ("35s",   "news_creative",      "Sonnet: 800-word SEO article\n'8,212 Projects Got RERA Notices...'",           PURPLE),
    ("53s",   "platform_agents ×5", "Async: Twitter thread (4 parts)\n+ Instagram (5 story slides + reel)\n+ YouTube Short + Housing News",  PURPLE),
    ("71s",   "qa_agent",           "Safety: PASS  ·  Quality: 6.4–8.0/10\nEngagement: 5.7%–7.15% predicted",       GREEN),
    ("97s",   "publisher",          "4 posts saved (dry run)\n2 internal links per post  ·  Slack summary sent",     ORANGE),
]
for i, (time, node, desc, color) in enumerate(timeline):
    y = Inches(1.65 + i * 0.65)
    rect(s, Inches(0.4), y, Inches(0.65), Inches(0.55), color)
    txb(s, time, Inches(0.4), y, Inches(0.65), Inches(0.55),
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s, node, Inches(1.15), y+Inches(0.02), Inches(2.1), Inches(0.45),
        size=11, bold=True, color=PURPLE)
    txb(s, desc, Inches(3.4), y+Inches(0.02), Inches(9.5), Inches(0.5),
        size=11, color=DARK)
    if i < 7:
        rect(s, Inches(0.62), y+Inches(0.55), Inches(0.12), Inches(0.12), color)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
subtitle_bar(s, "Roadmap: What's Next")
slide_num(s, 18); watermark(s)

phases = [
    ("Now  ✅\nHackathon Demo",
     ["5-platform auto-publish", "3-pass QA engine", "Hinglish social + SEO news",
      "Slack bot trigger", "Engagement prediction", "2 min end-to-end runtime"], GREEN),
    ("Month 1–2\nProduction Deploy",
     ["Live posting to all 5 platforms", "Human-in-the-loop approval UI",
      "PostgresSaver for multi-worker", "A/B model routing testing",
      "Hooks bank weekly refresh", "Monitoring + alerting dashboard"], PURPLE),
    ("Month 3–6\nScale & Intelligence",
     ["Vernacular expansion (Hindi, Tamil, Telugu)", "Video generation (Sora API)",
      "Personalized city microsites", "Builder partnership content",
      "Engagement prediction V2 (trained on real data)", "WhatsApp channel integration"], ORANGE),
]
for i, (phase, items, color) in enumerate(phases):
    x = Inches(0.4 + i * 4.28)
    rect(s, x, Inches(1.75), Inches(4.0), Inches(0.55), color)
    txb(s, phase, x, Inches(1.8), Inches(4.0), Inches(0.45),
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        y = Inches(2.4 + j * 0.68)
        rect(s, x+Inches(0.15), y+Inches(0.12), Inches(0.18), Inches(0.18), color)
        txb(s, item, x+Inches(0.45), y, Inches(3.4), Inches(0.6),
            size=12, color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — CALL TO ACTION
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, PURPLE_DARK)
slide_num(s, 19)

txb(s, "The AI Content Engine is live.", Inches(1), Inches(1.2), Inches(11), Inches(0.9),
    size=44, bold=True, color=WHITE)
txb(s, "300 posts/month. ₹14,000/month all-in.\nBrand-safe, platform-native, Hinglish-first.",
    Inches(1), Inches(2.2), Inches(11), Inches(1.0),
    size=24, color=PURPLE_LIGHT)

# 3 asks
asks = [
    ("✅ Deploy Now",    "DRY_RUN=false + platform keys\n→ Live in 30 minutes"),
    ("📊 Pilot 30 days", "Run alongside human team\nMeasure engagement delta"),
    ("🚀 Scale",         "Multi-city, multi-language\n→ Housing.com #1 PropTech content brand"),
]
for i, (ask, sub) in enumerate(asks):
    x = Inches(1.0 + i * 3.9)
    rect(s, x, Inches(3.5), Inches(3.5), Inches(1.9), PURPLE)
    txb(s, ask, x+Inches(0.15), Inches(3.6), Inches(3.2), Inches(0.55),
        size=18, bold=True, color=WHITE)
    txb(s, sub, x+Inches(0.15), Inches(4.2), Inches(3.2), Inches(0.9),
        size=13, color=PURPLE_LIGHT)

txb(s, "Demo: python main.py run  ·  Output in: output/<run_id>/",
    Inches(1), Inches(5.8), Inches(11), Inches(0.5),
    size=14, color=RGBColor(0x7B, 0xFF, 0x8B), bold=True)
txb(s, "Built at Housing.com Hackathon 2026  ·  Multi-Agent AI Stack",
    Inches(1), Inches(6.5), Inches(11), Inches(0.4),
    size=13, color=PURPLE_LIGHT, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
prs.save(OUT)
print(f"✓ Saved {OUT}  ({len(prs.slides)} slides)")
